import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import tkinter as tk
from tkinter import filedialog, messagebox


def adjust_picture_border(file_paths, weight, log_text):
    black_color = RGBColor(0, 0, 0)
    EMU = 918400  # English Metric Units (1 inch = 914400 EMUs)

    for pptx_path in file_paths:
        log_text.insert(tk.END, f"Processing {pptx_path}...\n")
        log_text.update()

        try:
            prs = Presentation(pptx_path)
        except Exception as e:
            log_text.insert(tk.END, f"❌ Error opening {pptx_path}: {e}\n")
            log_text.update()
            continue

        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for slide in prs.slides:
            shapes_list = [
                s for s in slide.shapes
                if s.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART)
            ]
            total_shapes = len(shapes_list)

            large_picture_positioned = False
            chart_image_positioned = False

            # Predefined positions
            large_image_left, large_image_top = int(3.19 * EMU), int(1.97 * EMU)
            small_image_left, small_image_top = int(13.87 * EMU), int(1.96 * EMU)
            chart_image_left, chart_image_top = int(2.4 * EMU), int(2.45 * EMU)

            # Single (centered) image/chart standard size
            single_width = int(10.5 * EMU)
            single_height = int(6 * EMU)

            for shape in shapes_list:
                try:
                    # 🟦 Case 1: Single image/chart slide
                    if total_shapes == 1:
                        shape.width = single_width
                        shape.height = single_height
                        shape.left = int((slide_width - shape.width) / 2)
                        shape.top = int((slide_height - shape.height) / 2)

                    # 🟨 Case 2: Multi-image slide (your previous logic)
                    else:
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            original_width, original_height = shape.width, shape.height
                            if original_width >= 2.74 * EMU and original_height >= 2.45 * EMU:
                                scale_factor_width, scale_factor_height = 2.06, 2.04
                                if not large_picture_positioned:
                                    shape.left, shape.top = large_image_left, large_image_top
                                    large_picture_positioned = True
                            else:
                                scale_factor_width, scale_factor_height = 1.16, 1.16
                                shape.left, shape.top = (
                                    small_image_left if large_picture_positioned else 0,
                                    small_image_top if large_picture_positioned else 0,
                                )

                            shape.width = int(original_width * scale_factor_width)
                            shape.height = int(original_height * scale_factor_height)

                        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                            if not chart_image_positioned:
                                shape.width, shape.height = int(10.73 * EMU), int(7.08 * EMU)
                                shape.left, shape.top = chart_image_left, chart_image_top
                                chart_image_positioned = True

                    # 🟥 Apply black border
                    if hasattr(shape, "line"):
                        shape.line.color.rgb = black_color
                        shape.line.width = Pt(float(weight))

                except Exception as e:
                    log_text.insert(tk.END, f"⚠️ Error processing shape: {e}\n")

        try:
            prs.save(pptx_path)
            log_text.insert(tk.END, f"✅ Processed successfully: {pptx_path}\n\n")
            log_text.update()
        except PermissionError:
            log_text.insert(tk.END, f"⚠️ File open in PowerPoint. Please close and retry.\n")
        except Exception as e:
            log_text.insert(tk.END, f"❌ Error saving {pptx_path}: {e}\n")

    messagebox.showinfo("Processing Complete", "All PowerPoint files processed successfully!")


def select_pptx_files(log_text):
    file_paths = filedialog.askopenfilenames(filetypes=[("PowerPoint Files", "*.pptx")])
    return file_paths


def start_processing(log_text, weight_entry):
    file_paths = select_pptx_files(log_text)
    if not file_paths:
        messagebox.showwarning("No Files Selected", "Please select at least one PowerPoint file.")
        return

    try:
        weight = float(weight_entry.get())
        adjust_picture_border(file_paths, weight, log_text)
    except ValueError:
        messagebox.showerror("Invalid Input", "Please enter a valid numeric border width.")


def main():
    root = tk.Tk()
    root.title("PowerPoint Image & Chart Resizer")
    root.geometry("550x400")

    tk.Label(root, text="Border Width (pts):").grid(row=0, column=0, padx=5, pady=5, sticky="e")
    weight_entry = tk.Entry(root)
    weight_entry.grid(row=0, column=1, padx=5, pady=5, sticky="w")
    weight_entry.insert(0, "1.5")

    log_text = tk.Text(root, width=65, height=18, wrap="word")
    log_text.grid(row=1, column=0, columnspan=2, padx=8, pady=5)

    select_button = tk.Button(
        root,
        text="Select & Process PPTX",
        bg="#0078D7",
        fg="white",
        font=("Arial", 10, "bold"),
        relief="groove",
        command=lambda: start_processing(log_text, weight_entry)
    )
    select_button.grid(row=2, column=0, columnspan=2, pady=10)

    root.mainloop()


if __name__ == "__main__":
    main()
